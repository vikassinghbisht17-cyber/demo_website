from pptx import Presentation

prs = Presentation(r'e:\pssgway_website__RECREATE\pssgway_docs\Data Case studies v1.pptx')
with open('extracted_slides.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        texts = []
        for s in slide.shapes:
            if hasattr(s, 'text'):
                texts.append(s.text.replace('\n', ' '))
        f.write(f'Slide {i+1}: ' + ' '.join(texts) + '\n')
