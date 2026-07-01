import os
import sys
import subprocess
import site

sys.path.append(site.getusersitepackages())

def install_deps():
    try:
        import pypdf
        import pptx
        print("Required dependencies already installed.")
    except ImportError:
        print("Installing dependencies pypdf and python-pptx...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--user", "pypdf", "python-pptx"])

def extract_pdf(pdf_path):
    print(f"\n=================== EXTRACTING PDF: {pdf_path} ===================")
    if not os.path.exists(pdf_path):
        print(f"File not found: {pdf_path}")
        return ""
    
    import pypdf
    reader = pypdf.PdfReader(pdf_path)
    text_content = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        text_content.append(f"--- PAGE {i+1} ---\n{text}\n")
    return "\n".join(text_content)

def extract_pptx(pptx_path):
    print(f"\n=================== EXTRACTING PPTX: {pptx_path} ===================")
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return ""
    
    from pptx import Presentation
    prs = Presentation(pptx_path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_text.append(f"--- SLIDE {i+1} ---")
        
        # Extract title if present
        if slide.shapes.title:
            slide_text.append(f"TITLE: {slide.shapes.title.text}")
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Avoid duplicating title
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                slide_text.append(shape.text.strip())
        
        # Extract notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"NOTES: {notes}")
                
        text_content.append("\n".join(slide_text) + "\n")
    return "\n".join(text_content)

if __name__ == "__main__":
    install_deps()
    
    files_to_extract = [
        "Passageway_Tech_Deck_Final.pdf",
        "Passageway-Data-Story.pdf",
        "Passageway-NBFC-Data-Story.pptx"
    ]
    
    for f in files_to_extract:
        base_dir = "e:\\pssgway_website__RECREATE"
        full_path = os.path.join(base_dir, f)
        
        if f.endswith(".pdf"):
            text = extract_pdf(full_path)
        else:
            text = extract_pptx(full_path)
            
        out_name = f.replace(".pdf", "_txt.txt").replace(".pptx", "_txt.txt")
        out_path = os.path.join(base_dir, out_name)
        with open(out_path, "w", encoding="utf-8") as out_file:
            out_file.write(text)
        print(f"Saved extracted text to: {out_path}")
